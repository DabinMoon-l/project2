"""
PPTX 퀴즈 생성 Cloud Run 서비스
- PPTX 파일에서 텍스트 추출
- Gemini API로 퀴즈 생성
- Firestore에 결과 저장
"""

import os
import json
import tempfile
import re
import subprocess
import uuid
from datetime import datetime
from flask import Flask, request, jsonify
from pptx import Presentation
from pptx.util import Inches
import google.generativeai as genai
import firebase_admin
from firebase_admin import credentials, firestore, storage, auth as fb_auth
from flask import send_file

# Flask 앱 초기화
app = Flask(__name__)

# Firebase 초기화 (Cloud Run 환경에서는 기본 서비스 계정 사용)
if not firebase_admin._apps:
    firebase_admin.initialize_app()

db = firestore.client()

# Gemini API 설정
GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)


def extract_text_from_pptx(file_path: str) -> list[dict]:
    """
    PPTX 파일에서 슬라이드별 텍스트 추출

    Returns:
        list of { 'slide_num': int, 'title': str, 'content': str }
    """
    prs = Presentation(file_path)
    slides_data = []

    # 반복되는 푸터/페이지번호 패턴 감지용
    footer_candidates = {}

    for slide_num, slide in enumerate(prs.slides, 1):
        title = ""
        content_parts = []

        for shape in slide.shapes:
            # 제목 추출
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    title = shape.text.strip()
                else:
                    # 일반 텍스트 추출
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            content_parts.append(text)

            # 표(Table) 텍스트 추출
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        content_parts.append(" | ".join(row_texts))

        # 푸터 후보 추적 (여러 슬라이드에서 반복되는 텍스트)
        for text in content_parts:
            if len(text) < 50:  # 짧은 텍스트만 푸터 후보로
                footer_candidates[text] = footer_candidates.get(text, 0) + 1

        content = "\n".join(content_parts)

        if title or content:
            slides_data.append({
                'slide_num': slide_num,
                'title': title,
                'content': content
            })

    # 반복되는 푸터 제거 (전체 슬라이드의 50% 이상에서 반복되는 텍스트)
    total_slides = len(slides_data)
    footers_to_remove = {
        text for text, count in footer_candidates.items()
        if count >= total_slides * 0.5 and count > 2
    }

    # 페이지 번호 패턴 제거
    page_num_pattern = re.compile(r'^\d+$|^페이지\s*\d+|^Page\s*\d+|^\d+\s*/\s*\d+$', re.IGNORECASE)

    for slide in slides_data:
        lines = slide['content'].split('\n')
        filtered_lines = [
            line for line in lines
            if line not in footers_to_remove and not page_num_pattern.match(line.strip())
        ]
        slide['content'] = "\n".join(filtered_lines)

    return slides_data


def chunk_slides(slides_data: list[dict], chunk_size: int = 3) -> list[str]:
    """
    슬라이드를 chunk_size개씩 묶어서 텍스트 청크 생성
    """
    chunks = []

    for i in range(0, len(slides_data), chunk_size):
        chunk_slides = slides_data[i:i + chunk_size]
        chunk_text = ""

        for slide in chunk_slides:
            chunk_text += f"\n[슬라이드 {slide['slide_num']}]\n"
            if slide['title']:
                chunk_text += f"제목: {slide['title']}\n"
            if slide['content']:
                chunk_text += f"{slide['content']}\n"

        if chunk_text.strip():
            chunks.append(chunk_text.strip())

    return chunks


def generate_quiz_from_chunk(chunk_text: str, difficulty: str = 'medium', keywords: list[str] = None) -> list[dict]:
    """
    텍스트 청크에서 Gemini API로 퀴즈 생성
    """
    difficulty_desc = {
        'easy': '기본 개념을 확인하는 쉬운',
        'medium': '이해도를 테스트하는 보통 난이도의',
        'hard': '응용력을 요구하는 어려운'
    }

    # 키워드가 있으면 프롬프트에 포함
    keyword_instruction = ""
    if keywords and len(keywords) > 0:
        keyword_list = ", ".join(keywords)
        keyword_instruction = f"\n4. 다음 키워드들을 중심으로 문제를 출제해주세요: {keyword_list}"

    prompt = f"""다음 학습 자료를 바탕으로 {difficulty_desc.get(difficulty, '보통 난이도의')} 객관식 문제를 2~3개 생성해주세요.

[학습 자료]
{chunk_text}

[요구사항]
1. 각 문제는 학습 내용의 핵심 개념을 테스트해야 합니다.
2. 선지는 5개로 구성하고, 오답도 그럴듯하게 만들어주세요.
3. 반드시 아래 JSON 형식으로만 응답해주세요.{keyword_instruction}

[JSON 형식]
```json
[
  {{
    "question": "문제 내용",
    "choices": ["선지1", "선지2", "선지3", "선지4", "선지5"],
    "answer_index": 0,
    "explanation": "정답 해설",
    "topic": "관련 주제"
  }}
]
```
"""

    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content(prompt)

        # JSON 추출
        response_text = response.text

        # ```json ... ``` 블록 추출
        json_match = re.search(r'```json\s*([\s\S]*?)\s*```', response_text)
        if json_match:
            json_str = json_match.group(1)
        else:
            # JSON 배열 직접 찾기
            json_match = re.search(r'\[[\s\S]*\]', response_text)
            if json_match:
                json_str = json_match.group(0)
            else:
                print(f"JSON을 찾을 수 없음: {response_text[:200]}")
                return []

        questions = json.loads(json_str)

        # 유효성 검증
        validated_questions = []
        for q in questions:
            if all(key in q for key in ['question', 'choices', 'answer_index', 'explanation']):
                if len(q['choices']) >= 4 and 0 <= q['answer_index'] < len(q['choices']):
                    validated_questions.append({
                        'question': q['question'],
                        'choices': q['choices'][:5],  # 최대 5개
                        'answer_index': q['answer_index'],
                        'explanation': q['explanation'],
                        'topic': q.get('topic', '')
                    })

        return validated_questions

    except Exception as e:
        print(f"Gemini API 오류: {e}")
        return []


def update_job_status(job_id: str, status: str, progress: int = 0, **kwargs):
    """Firestore 작업 상태 업데이트"""
    update_data = {
        'status': status,
        'progress': progress,
        'updatedAt': firestore.SERVER_TIMESTAMP
    }
    update_data.update(kwargs)

    db.collection('quizJobs').document(job_id).update(update_data)


def verify_firebase_token(req):
    """Authorization 헤더에서 Firebase ID 토큰 검증"""
    auth_header = req.headers.get('Authorization', '')
    if not auth_header.startswith('Bearer '):
        return None
    token = auth_header[7:]
    try:
        return fb_auth.verify_id_token(token)
    except Exception:
        return None


@app.route('/convert-pdf', methods=['POST', 'OPTIONS'])
def convert_pdf():
    """
    PPT → PDF 직접 변환 (클라이언트 직접 호출)
    PPTX 바이너리를 multipart로 받아 PDF 바이너리 반환
    Firebase Auth 토큰으로 인증
    """
    # CORS preflight
    if request.method == 'OPTIONS':
        resp = app.make_default_options_response()
        resp.headers['Access-Control-Allow-Origin'] = '*'
        resp.headers['Access-Control-Allow-Methods'] = 'POST, OPTIONS'
        resp.headers['Access-Control-Allow-Headers'] = 'Authorization, Content-Type'
        resp.headers['Access-Control-Max-Age'] = '3600'
        return resp

    # Firebase Auth 검증
    decoded = verify_firebase_token(request)
    if not decoded:
        return jsonify({'error': '인증 실패'}), 401

    try:
        # multipart 파일 수신
        if 'file' not in request.files:
            return jsonify({'error': 'file이 필요합니다.'}), 400

        pptx_file = request.files['file']
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            pptx_file.save(tmp.name)
            pptx_path = tmp.name

        try:
            # LibreOffice로 PDF 변환
            output_dir = tempfile.mkdtemp()
            subprocess.run(
                ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
                check=True,
                timeout=120,
                capture_output=True,
                text=True,
            )

            pdf_filename = os.path.basename(pptx_path).replace('.pptx', '.pdf')
            pdf_path = os.path.join(output_dir, pdf_filename)

            if not os.path.exists(pdf_path):
                return jsonify({'error': 'PDF 변환 결과 파일을 찾을 수 없습니다.'}), 500

            # PDF 바이너리 직접 반환
            resp = send_file(
                pdf_path,
                mimetype='application/pdf',
                as_attachment=False,
            )
            resp.headers['Access-Control-Allow-Origin'] = '*'

            # 임시 PPTX 파일 정리 (send_file 후 정리는 Flask가 처리)
            # output_dir 정리는 after_request로 처리
            @resp.call_on_close
            def cleanup():
                try:
                    if os.path.exists(pdf_path):
                        os.unlink(pdf_path)
                    if os.path.exists(output_dir):
                        os.rmdir(output_dir)
                except Exception:
                    pass

            return resp

        finally:
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)

    except subprocess.TimeoutExpired:
        resp = jsonify({'error': 'PDF 변환 시간 초과 (120초)'})
        resp.headers['Access-Control-Allow-Origin'] = '*'
        return resp, 504
    except subprocess.CalledProcessError as e:
        print(f"LibreOffice 변환 오류: {e.stderr}")
        resp = jsonify({'error': f'LibreOffice 변환 실패: {e.stderr}'})
        resp.headers['Access-Control-Allow-Origin'] = '*'
        return resp, 500
    except Exception as e:
        print(f"PDF 변환 오류: {e}")
        resp = jsonify({'error': str(e)})
        resp.headers['Access-Control-Allow-Origin'] = '*'
        return resp, 500


@app.route('/convert-to-pdf-direct', methods=['POST'])
def convert_to_pdf_direct():
    """
    PPT → PDF 변환 (JSON base64 I/O, Cloud Functions에서 호출)
    Request: { "pptxBase64": "..." }
    Response: { "success": true, "pdfBase64": "..." }
    """
    import base64

    data = request.get_json(force=True)
    pptx_base64 = data.get('pptxBase64')
    if not pptx_base64:
        return jsonify({'error': 'pptxBase64가 필요합니다.'}), 400

    try:
        pptx_bytes = base64.b64decode(pptx_base64)

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(pptx_bytes)
            pptx_path = tmp.name

        try:
            output_dir = tempfile.mkdtemp()
            subprocess.run(
                ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
                check=True,
                timeout=120,
                capture_output=True,
                text=True,
            )

            pdf_filename = os.path.basename(pptx_path).replace('.pptx', '.pdf')
            pdf_path = os.path.join(output_dir, pdf_filename)

            if not os.path.exists(pdf_path):
                return jsonify({'error': 'PDF 변환 결과 파일을 찾을 수 없습니다.'}), 500

            with open(pdf_path, 'rb') as f:
                pdf_base64 = base64.b64encode(f.read()).decode('utf-8')

            # 임시 파일 정리
            os.unlink(pdf_path)
            os.rmdir(output_dir)

            return jsonify({'success': True, 'pdfBase64': pdf_base64})

        finally:
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)

    except subprocess.TimeoutExpired:
        return jsonify({'error': 'PDF 변환 시간 초과 (120초)'}), 504
    except subprocess.CalledProcessError as e:
        print(f"LibreOffice 변환 오류: {e.stderr}")
        return jsonify({'error': f'LibreOffice 변환 실패: {e.stderr}'}), 500
    except Exception as e:
        print(f"PDF 변환 오류: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/process-pptx', methods=['POST'])
def process_pptx():
    """
    PPTX 처리 엔드포인트

    Request Body:
    {
        "jobId": "작업 ID",
        "storagePath": "pptx-uploads/uid/filename.pptx",
        "userId": "사용자 UID",
        "folderName": "퀴즈 이름",
        "difficulty": "easy|medium|hard",
        "questionCount": 10,
        "tags": ["태그1", "태그2"]
    }
    """
    try:
        data = request.get_json()

        job_id = data.get('jobId')
        storage_path = data.get('storagePath')
        user_id = data.get('userId')
        folder_name = data.get('folderName', '퀴즈')
        difficulty = data.get('difficulty', 'medium')
        question_count = data.get('questionCount', 10)
        tags = data.get('tags', [])
        keywords = data.get('keywords', [])  # 추출된 키워드

        if not all([job_id, storage_path, user_id]):
            return jsonify({'error': 'Missing required fields'}), 400

        # 상태 업데이트: 처리 시작
        update_job_status(job_id, 'processing', 5, message='PPTX 다운로드 중...')

        # Storage에서 PPTX 다운로드
        bucket_name = os.environ.get('FIREBASE_STORAGE_BUCKET', 'project2-7a317.firebasestorage.app')
        bucket = storage.bucket(bucket_name)
        blob = bucket.blob(storage_path)

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            blob.download_to_filename(tmp_file.name)
            tmp_path = tmp_file.name

        try:
            # 상태 업데이트: 텍스트 추출
            update_job_status(job_id, 'processing', 15, message='슬라이드 텍스트 추출 중...')

            # PPTX 텍스트 추출 (전체 슬라이드)
            slides_data = extract_text_from_pptx(tmp_path)

            if not slides_data:
                update_job_status(job_id, 'failed', 0, error='슬라이드에서 텍스트를 찾을 수 없습니다.')
                return jsonify({'error': 'No text found in PPTX'}), 400

            # 슬라이드 청킹
            chunks = chunk_slides(slides_data, chunk_size=3)

            keyword_info = f' (키워드 {len(keywords)}개)' if keywords else ''
            update_job_status(
                job_id, 'processing', 25,
                message=f'{len(slides_data)}개 슬라이드에서 {len(chunks)}개 청크 생성됨{keyword_info}'
            )

            # 퀴즈 생성
            all_questions = []
            for i, chunk in enumerate(chunks):
                progress = 25 + int((i / len(chunks)) * 60)
                update_job_status(
                    job_id, 'processing', progress,
                    message=f'퀴즈 생성 중... ({i + 1}/{len(chunks)})'
                )

                questions = generate_quiz_from_chunk(chunk, difficulty, keywords)
                all_questions.extend(questions)

                # 요청 개수 도달 시 중단
                if len(all_questions) >= question_count:
                    break

            # 문제 수 조정
            all_questions = all_questions[:question_count]

            if not all_questions:
                update_job_status(job_id, 'failed', 0, error='문제 생성에 실패했습니다.')
                return jsonify({'error': 'Failed to generate questions'}), 500

            # 상태 업데이트: Firestore 저장
            update_job_status(job_id, 'processing', 90, message='퀴즈 저장 중...')

            # Firestore에 퀴즈 저장 (기존 learning 타입 구조)
            quiz_ref = db.collection('quizzes').document()
            quiz_data = {
                'title': folder_name,
                'type': 'learning',
                'tags': tags,
                'questions': [
                    {
                        'id': f'q{i + 1}',
                        'type': 'multiple',
                        'text': q['question'],
                        'choices': q['choices'],
                        'answer': q['answer_index'],
                        'explanation': q['explanation'],
                        'topic': q.get('topic', '')
                    }
                    for i, q in enumerate(all_questions)
                ],
                'creatorId': user_id,
                'isPublic': False,
                'totalQuestions': len(all_questions),
                'sourceType': 'pptx',
                'slideCount': len(slides_data),
                'keywords': keywords,  # 사용된 키워드 저장
                'createdAt': firestore.SERVER_TIMESTAMP,
                'updatedAt': firestore.SERVER_TIMESTAMP
            }

            quiz_ref.set(quiz_data)

            # Storage에서 PPTX 파일 삭제 (비용 절약)
            try:
                blob.delete()
            except Exception as e:
                print(f"Storage 파일 삭제 실패: {e}")

            # 상태 업데이트: 완료
            update_job_status(
                job_id, 'completed', 100,
                message='퀴즈 생성 완료!',
                quizId=quiz_ref.id,
                questionCount=len(all_questions)
            )

            return jsonify({
                'success': True,
                'quizId': quiz_ref.id,
                'questionCount': len(all_questions)
            })

        finally:
            # 임시 파일 삭제
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

    except Exception as e:
        print(f"처리 오류: {e}")
        if job_id:
            update_job_status(job_id, 'failed', 0, error=str(e))
        return jsonify({'error': str(e)}), 500


@app.route('/render-pdf', methods=['POST', 'OPTIONS'])
def render_pdf():
    """
    HTML → PDF 변환 (Playwright/Chromium)
    Request: POST { "html": "<html>...</html>" }
    Response: PDF 바이너리 (application/pdf)
    Firebase Auth 토큰으로 인증
    """
    # CORS preflight
    if request.method == 'OPTIONS':
        resp = app.make_default_options_response()
        resp.headers['Access-Control-Allow-Origin'] = '*'
        resp.headers['Access-Control-Allow-Methods'] = 'POST, OPTIONS'
        resp.headers['Access-Control-Allow-Headers'] = 'Authorization, Content-Type'
        resp.headers['Access-Control-Max-Age'] = '3600'
        return resp

    # Firebase Auth 검증
    decoded = verify_firebase_token(request)
    if not decoded:
        resp = jsonify({'error': '인증 실패'})
        resp.headers['Access-Control-Allow-Origin'] = '*'
        return resp, 401

    try:
        data = request.get_json(force=True)
        html_content = data.get('html', '')
        if not html_content:
            resp = jsonify({'error': 'html 필드가 필요합니다.'})
            resp.headers['Access-Control-Allow-Origin'] = '*'
            return resp, 400

        from playwright.sync_api import sync_playwright

        with sync_playwright() as p:
            browser = p.chromium.launch(args=['--no-sandbox', '--disable-dev-shm-usage'])
            page = browser.new_page()
            page.set_content(html_content, wait_until='networkidle')
            pdf_bytes = page.pdf(
                format='A4',
                margin={'top': '0', 'right': '0', 'bottom': '0', 'left': '0'},
                print_background=True,
            )
            browser.close()

        # PDF 바이너리 직접 반환
        import io
        resp = send_file(
            io.BytesIO(pdf_bytes),
            mimetype='application/pdf',
            as_attachment=False,
        )
        resp.headers['Access-Control-Allow-Origin'] = '*'
        return resp

    except Exception as e:
        print(f"HTML→PDF 변환 오류: {e}")
        resp = jsonify({'error': str(e)})
        resp.headers['Access-Control-Allow-Origin'] = '*'
        return resp, 500


@app.route('/health', methods=['GET'])
def health_check():
    """헬스 체크 엔드포인트"""
    return jsonify({'status': 'healthy'})


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 8080))
    app.run(host='0.0.0.0', port=port, debug=False)
